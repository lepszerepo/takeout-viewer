"""mbox parser v2 — streams large mailboxes one message at a time.

Per message we capture:
- MIME-decoded Subject, From, To, Cc, Bcc, Reply-To, Date, Message-ID, In-Reply-To, References
- Plain-text body (preferring text/plain part), with HTML→text fallback
- Raw HTML body (text/html part) for fidelity in UI preview
- Attachment names + sizes (no payloads — keep DB lean)
- Selected raw headers (List-Id, X-Gmail-Labels, Authentication-Results)
- Folder hint inferred from filename (Inbox/Trash/Quarantine/Deleted/Spam)

The emitted NormalizedEvent carries display fields (title, description, etc.).
Full body + headers + attachments are also passed through `extra` so the
importer can persist them into the `mail_messages` table.
"""
from __future__ import annotations

import hashlib
import re
from email import policy
from email.header import decode_header, make_header
from email.parser import BytesParser
from email.utils import parsedate_to_datetime, getaddresses
from pathlib import Path
from typing import Iterator, Optional

from ...config import settings

from ..normalizer import NormalizedEvent

_BODY_SNIPPET_CHARS = 300
_MAX_BODY_BYTES = 4 * 1024 * 1024  # cap stored body per message at 4 MB
_MAX_MSG_BYTES = 64 * 1024 * 1024  # cap per-message buffer at 64 MB
_MAX_ATTACH_SAVE = 100 * 1024 * 1024  # don't save attachments larger than 100 MB


def _attachment_path(sha256_hex: str) -> Path:
    base = settings.data_dir / "attachments" / sha256_hex[:2]
    base.mkdir(parents=True, exist_ok=True)
    return base / sha256_hex[2:]


def _save_attachment(payload: bytes) -> str | None:
    """Write payload to a content-addressed file. Returns sha256 hex."""
    if not payload:
        return None
    sha = hashlib.sha256(payload).hexdigest()
    path = _attachment_path(sha)
    if not path.exists():
        try:
            path.write_bytes(payload)
        except Exception:
            return None
    return sha


def _decode(s: Optional[str]) -> str:
    if not s:
        return ""
    try:
        return str(make_header(decode_header(s))).strip()
    except Exception:
        try:
            return str(s).strip()
        except Exception:
            return ""


def _hv(msg, name: str) -> str:
    v = msg.get(name)
    return _decode(v if isinstance(v, str) else (str(v) if v is not None else ""))


def _hv_raw(msg, name: str) -> str:
    v = msg.get(name)
    if v is None:
        return ""
    try:
        return str(v).strip()
    except Exception:
        return ""


def _parse_date(value: Optional[str]):
    if not value:
        return None
    try:
        dt = parsedate_to_datetime(value)
    except (TypeError, ValueError):
        return None
    if dt is None:
        return None
    if dt.tzinfo is not None:
        from datetime import timezone
        dt = dt.astimezone(timezone.utc).replace(tzinfo=None)
    return dt


def _addr_list(value: str) -> list[str]:
    if not value:
        return []
    try:
        parsed = getaddresses([value])
    except Exception:
        return []
    return [a[1].strip() for a in parsed if a[1]]


def _addrs_with_names(value: str) -> list[dict]:
    if not value:
        return []
    try:
        parsed = getaddresses([value])
    except Exception:
        return []
    out = []
    for name, email in parsed:
        if not email:
            continue
        out.append({"name": _decode(name), "email": email.strip()})
    return out


def _html_to_text(html: str) -> str:
    """Strip HTML tags for body_text fallback."""
    try:
        from bs4 import BeautifulSoup
        return BeautifulSoup(html, "lxml").get_text("\n", strip=True)
    except Exception:
        return ""


def _references(value: str) -> list[str]:
    """Parse References / In-Reply-To header: whitespace-separated <id> tokens."""
    if not value:
        return []
    out = []
    for tok in value.split():
        tok = tok.strip("<>").strip()
        if tok:
            out.append(tok)
    return out


_MAX_ATTACH_TEXT = 200_000  # cap text per attachment


def _attachment_text(payload: bytes, content_type: str, filename: str) -> str:
    """Try to extract text from an attachment payload. Best-effort, never raises."""
    import io
    import tempfile
    if not payload:
        return ""
    ct = (content_type or "").lower()
    name_lc = (filename or "").lower()
    ext = name_lc.rsplit(".", 1)[-1] if "." in name_lc else ""
    try:
        if "pdf" in ct or ext == "pdf":
            from pdfminer.high_level import extract_text
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tf:
                tf.write(payload); tf.flush()
                return (extract_text(tf.name) or "")[:_MAX_ATTACH_TEXT]
        if "wordprocessingml" in ct or ext == "docx":
            from docx import Document
            doc = Document(io.BytesIO(payload))
            return ("\n".join(p.text for p in doc.paragraphs))[:_MAX_ATTACH_TEXT]
        if ext in {"xlsx", "xlsm"} or "spreadsheetml" in ct:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(payload), read_only=True, data_only=True)
            buf: list[str] = []
            for sheet in wb.worksheets:
                buf.append(f"# {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        buf.append("\t".join(cells))
            return ("\n".join(buf))[:_MAX_ATTACH_TEXT]
        if "presentationml" in ct or ext == "pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(payload))
            out: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            txt = "".join(r.text for r in p.runs).strip()
                            if txt:
                                out.append(txt)
            return ("\n".join(out))[:_MAX_ATTACH_TEXT]
        if ct.startswith("text/") or ext in {"txt", "csv", "md", "json", "html", "xml", "log", "tsv"}:
            return payload.decode("utf-8", errors="replace")[:_MAX_ATTACH_TEXT]
        if ct.startswith("image/") or ext in {"jpg", "jpeg", "png", "heic", "heif", "tif", "tiff", "bmp", "webp"}:
            import pytesseract
            from PIL import Image
            try:
                from pillow_heif import register_heif_opener
                register_heif_opener()
            except Exception:
                pass
            try:
                img = Image.open(io.BytesIO(payload))
                return pytesseract.image_to_string(img, lang="pol+eng")[:_MAX_ATTACH_TEXT]
            except Exception:
                return ""
    except Exception:
        return ""
    return ""


def _extract_bodies(msg) -> tuple[str, str, list[dict]]:
    """Return (body_text, body_html, attachments).
    Attachments include extracted text where possible."""
    text_parts: list[str] = []
    html_parts: list[str] = []
    attachments: list[dict] = []

    try:
        walker = msg.walk() if msg.is_multipart() else [msg]
    except Exception:
        walker = [msg]

    for part in walker:
        try:
            ctype = part.get_content_type()
            disp = (part.get("Content-Disposition") or "").lower()
            filename = part.get_filename()
            if filename:
                filename = _decode(filename)

            if disp.startswith("attachment") or (filename and ctype not in ("text/plain", "text/html")):
                payload = part.get_payload(decode=True)
                size = len(payload) if isinstance(payload, (bytes, bytearray)) else 0
                att_text = ""
                sha = None
                if isinstance(payload, (bytes, bytearray)):
                    payload_b = bytes(payload)
                    if size < 50_000_000:
                        att_text = _attachment_text(payload_b, ctype, filename or "")
                    if size < _MAX_ATTACH_SAVE:
                        sha = _save_attachment(payload_b)
                attachments.append({
                    "name": filename or "(bez nazwy)",
                    "content_type": ctype,
                    "size": size,
                    "text": att_text or None,
                    "sha256": sha,
                })
                continue

            if ctype == "text/plain":
                payload = part.get_payload(decode=True)
                if isinstance(payload, (bytes, bytearray)):
                    charset = part.get_content_charset() or "utf-8"
                    text_parts.append(payload.decode(charset, errors="replace"))
            elif ctype == "text/html":
                payload = part.get_payload(decode=True)
                if isinstance(payload, (bytes, bytearray)):
                    charset = part.get_content_charset() or "utf-8"
                    html_parts.append(payload.decode(charset, errors="replace"))
        except Exception:
            continue

    body_text = "\n\n".join(p for p in text_parts if p).strip()
    body_html = "\n".join(p for p in html_parts if p).strip()
    if not body_text and body_html:
        body_text = _html_to_text(body_html)

    # Cap stored sizes
    if len(body_text.encode("utf-8")) > _MAX_BODY_BYTES:
        body_text = body_text.encode("utf-8")[:_MAX_BODY_BYTES].decode("utf-8", errors="replace")
    if len(body_html.encode("utf-8")) > _MAX_BODY_BYTES:
        body_html = body_html.encode("utf-8")[:_MAX_BODY_BYTES].decode("utf-8", errors="replace")

    return body_text, body_html, attachments


def _folder_hint(rel: str) -> str:
    rel_lc = rel.lower()
    name = rel_lc.rsplit("/", 1)[-1]
    if "trash" in name or "kosz" in name:
        return "Trash"
    if "quarantine" in name or "kwarantanna" in name:
        return "Quarantine"
    if "deleted" in name:
        return "Deleted"
    if "spam" in name:
        return "Spam"
    if "sent" in name or "wyslane" in name:
        return "Sent"
    if "inbox" in name or "all mail" in name or "cała poczta" in name:
        return "Inbox"
    return "Mail"


def _thread_id(refs: list[str], in_reply_to: list[str], msg_id: str) -> str:
    """Earliest ancestor message-id, or self if standalone."""
    if refs:
        return refs[0]
    if in_reply_to:
        return in_reply_to[0]
    return msg_id


def _build_event(msg_bytes: bytes, rel: str, folder: str, ordinal: int) -> Optional[NormalizedEvent]:
    if not msg_bytes:
        return None
    parser = BytesParser(policy=policy.compat32)
    try:
        msg = parser.parsebytes(msg_bytes)
    except Exception:
        return None

    subject = _hv(msg, "Subject")
    from_full = _addrs_with_names(_hv_raw(msg, "From"))
    to_full = _addrs_with_names(_hv_raw(msg, "To"))
    cc_full = _addrs_with_names(_hv_raw(msg, "Cc"))
    bcc_full = _addrs_with_names(_hv_raw(msg, "Bcc"))
    reply_to_full = _addrs_with_names(_hv_raw(msg, "Reply-To"))

    timestamp = _parse_date(_hv_raw(msg, "Date"))
    msg_id = _hv_raw(msg, "Message-ID").strip("<>") or _hv_raw(msg, "Message-Id").strip("<>")
    in_reply_to = _references(_hv_raw(msg, "In-Reply-To"))
    refs = _references(_hv_raw(msg, "References"))
    labels_raw = _hv_raw(msg, "X-Gmail-Labels")
    labels = [l.strip() for l in labels_raw.split(",") if l.strip()] if labels_raw else []

    body_text, body_html, attachments = _extract_bodies(msg)
    snippet = body_text[:_BODY_SNIPPET_CHARS] if body_text else None

    selected_headers = {
        h: _hv_raw(msg, h)
        for h in ("List-Id", "List-Unsubscribe", "Authentication-Results", "Return-Path", "X-Spam-Score")
        if _hv_raw(msg, h)
    }

    thread_id = _thread_id(refs, in_reply_to, msg_id or f"{rel}#{ordinal}")
    unique_marker = msg_id or f"{rel}#{ordinal}"

    from_emails = [a["email"] for a in from_full]
    to_emails = [a["email"] for a in to_full]

    ev = NormalizedEvent(
        source="mail",
        service="Gmail",
        category="mail",
        type="mail_message",
        title=subject or "(bez tematu)",
        description=snippet,
        timestamp=timestamp,
        people={
            "from": from_emails,
            "to": to_emails,
            "cc": [a["email"] for a in cc_full],
            "bcc": [a["email"] for a in bcc_full],
        },
        metadata={
            "folder": folder,
            "message_id": msg_id or None,
            "thread_id": thread_id,
            "labels": labels or None,
            "size_bytes": len(msg_bytes),
            "attachments_count": len(attachments),
            "has_html": bool(body_html),
        },
        raw_path=f"{rel}#{unique_marker}",
    )
    # Stash extra fields the importer will pick up for the mail_messages table
    ev.extra = {
        "mail_message_id": msg_id or None,
        "mail_in_reply_to": in_reply_to[0] if in_reply_to else None,
        "mail_references": refs,
        "mail_thread_id": thread_id,
        "mail_folder": folder,
        "mail_labels": labels,
        "mail_from": from_full,
        "mail_to": to_full,
        "mail_cc": cc_full,
        "mail_bcc": bcc_full,
        "mail_reply_to": reply_to_full,
        "mail_subject_raw": subject,
        "mail_body_text": body_text or None,
        "mail_body_html": body_html or None,
        "mail_attachments": attachments,
        "mail_headers": selected_headers,
        "mail_size_bytes": len(msg_bytes),
    }
    return ev


def mbox_parser(path: Path, rel: str) -> Iterator[NormalizedEvent]:
    folder = _folder_hint(rel)
    buf = bytearray()
    ordinal = 0
    with path.open("rb") as f:
        for line in f:
            if line.startswith(b"From ") and buf:
                ev = _build_event(bytes(buf), rel, folder, ordinal)
                ordinal += 1
                buf.clear()
                if ev is not None:
                    yield ev
            buf.extend(line)
            if len(buf) > _MAX_MSG_BYTES:
                ev = _build_event(bytes(buf), rel, folder, ordinal)
                ordinal += 1
                buf.clear()
                if ev is not None:
                    yield ev
        if buf:
            ev = _build_event(bytes(buf), rel, folder, ordinal)
            if ev is not None:
                yield ev


def register_parsers(register) -> None:
    register("mail", "mbox", mbox_parser)
