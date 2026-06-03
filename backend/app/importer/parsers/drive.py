"""Drive parser — extracts text from common file formats and stores it for search.

Supported text extraction:
- .pdf (pdfminer.six)
- .docx (python-docx)
- .xlsx (openpyxl) — cell values from all sheets
- .txt / .csv / .md / .json / .html / .xml (raw read with charset guess)

For each Drive file we emit ONE NormalizedEvent. The actual extracted text
+ metadata are persisted to the `drive_files` table via the importer.
"""
from __future__ import annotations

import json
from pathlib import Path
from typing import Iterator, Optional

from ..normalizer import NormalizedEvent
from .base import parse_iso

# Sidecar ".json-info" file gives us original ownership / sharing info
_INFO_SUFFIX = ".json-info"
_TEXT_LIMIT = 1_000_000  # cap stored extracted text at ~1 MB per doc


def _ocr_image(path: Path) -> str:
    import pytesseract
    from PIL import Image
    # HEIC support
    try:
        from pillow_heif import register_heif_opener
        register_heif_opener()
    except Exception:
        pass
    try:
        img = Image.open(str(path))
    except Exception:
        return ""
    try:
        return pytesseract.image_to_string(img, lang="pol+eng")
    except Exception:
        try:
            return pytesseract.image_to_string(img, lang="eng")
        except Exception:
            return ""


def _legacy_doc_text(path: Path) -> str:
    """Use antiword for .doc files."""
    import subprocess
    try:
        result = subprocess.run(
            ["antiword", "-w", "0", str(path)],
            capture_output=True, timeout=60, text=True,
        )
        return (result.stdout or "")
    except Exception:
        return ""


def _pptx_text(path: Path) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
    except Exception:
        return ""
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"--- Slajd {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in p.runs).strip()
                    if txt:
                        out.append(txt)
    return "\n".join(out)


def _extract_text(path: Path) -> tuple[str, str]:
    """Return (extracted_text, status). Status is 'ok' | 'unsupported' | 'error'."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            from pdfminer.high_level import extract_text
            text = (extract_text(str(path)) or "").strip()
            # If pdfminer returned very little, it's likely a scanned PDF — OCR
            if len(text) < 50 and path.stat().st_size > 5_000:
                ocr = _ocr_pdf(path)
                if ocr:
                    return ocr[:_TEXT_LIMIT], "ok"
            return text[:_TEXT_LIMIT], "ok"
        if ext == ".docx":
            from docx import Document
            doc = Document(str(path))
            return ("\n".join(p.text for p in doc.paragraphs) or "")[:_TEXT_LIMIT], "ok"
        if ext == ".doc":
            return _legacy_doc_text(path)[:_TEXT_LIMIT], "ok"
        if ext in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True, data_only=True)
            buf: list[str] = []
            for sheet in wb.worksheets:
                buf.append(f"# {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        buf.append("\t".join(cells))
                if sum(len(s) for s in buf) > _TEXT_LIMIT:
                    break
            return ("\n".join(buf))[:_TEXT_LIMIT], "ok"
        if ext in {".pptx", ".ppt"}:
            return _pptx_text(path)[:_TEXT_LIMIT], "ok"
        if ext in {".jpg", ".jpeg", ".png", ".heic", ".heif", ".tif", ".tiff", ".bmp", ".webp"}:
            return _ocr_image(path)[:_TEXT_LIMIT], "ok"
        if ext in {".txt", ".csv", ".md", ".json", ".html", ".xml", ".log", ".tsv", ".svg"}:
            with path.open("r", encoding="utf-8", errors="replace") as f:
                return f.read(_TEXT_LIMIT), "ok"
    except Exception:
        return "", "error"
    return "", "unsupported"


def _ocr_pdf(path: Path) -> str:
    """OCR a scanned PDF by rasterizing pages via poppler/pdftoppm then tesseract.

    Only handles up to 30 pages to keep import time bounded.
    """
    import subprocess
    import tempfile
    try:
        with tempfile.TemporaryDirectory() as td:
            subprocess.run(
                ["pdftoppm", "-r", "200", "-l", "30", str(path), f"{td}/page", "-png"],
                capture_output=True, timeout=300, check=False,
            )
            import pytesseract
            from PIL import Image
            parts: list[str] = []
            for p in sorted(Path(td).glob("page-*.png")):
                try:
                    img = Image.open(str(p))
                    parts.append(pytesseract.image_to_string(img, lang="pol+eng"))
                except Exception:
                    continue
            return "\n".join(parts)
    except Exception:
        return ""


def _load_info_sidecar(path: Path) -> Optional[dict]:
    sidecar = path.with_name(path.name + _INFO_SUFFIX)
    if not sidecar.exists():
        return None
    try:
        with sidecar.open("r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return None


def drive_parser(path: Path, rel: str) -> Iterator[NormalizedEvent]:
    # The scanner provides individual files; skip sidecar files themselves
    if rel.endswith(_INFO_SUFFIX):
        return
    name = path.name
    ext = path.suffix.lower().lstrip(".")
    size = path.stat().st_size

    info = _load_info_sidecar(path)
    timestamp = None
    if info:
        ts_raw = (
            info.get("modified_time")
            or info.get("modifiedTime")
            or info.get("created_time")
            or info.get("createdTime")
        )
        timestamp = parse_iso(ts_raw)

    text, status = _extract_text(path)

    description = None
    if text:
        description = text[:500]
    elif info and isinstance(info.get("description"), str):
        description = info["description"]

    owners = []
    if info:
        for k in ("owner", "owners", "shared_with", "permissions"):
            v = info.get(k)
            if isinstance(v, list):
                owners.extend([str(x) for x in v if x])
            elif isinstance(v, dict):
                owners.append(str(v.get("emailAddress") or v))
            elif isinstance(v, str):
                owners.append(v)

    yield NormalizedEvent(
        source="drive",
        service="Google Drive",
        category="file",
        type="drive_file",
        title=name,
        description=description,
        timestamp=timestamp,
        people=owners or None,
        metadata={
            "extension": ext,
            "size": size,
            "extraction_status": status,
            "info": info,
        },
        raw_path=rel,
        extra={
            "drive_extracted_text": text or None,
            "drive_extraction_status": status,
            "drive_info_json": info,
            "drive_size_bytes": size,
            "drive_extension": ext,
            "drive_file_name": name,
            "drive_relative_path": rel,
        },
    )


def register_parsers(register) -> None:
    register("drive", "file", drive_parser)
